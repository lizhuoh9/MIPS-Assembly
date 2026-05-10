"""
Build the MIPS Assembly teaching slides (30 slides, ~1 hour) on top of the
Hanyang BRL template. The template's slide master carries the BRL banner,
so every slide we add inherits it automatically.

Usage:  python _build_slides.py
Output: MIPS_Assembly_Lecture.pptx in this directory.
"""

from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = Path(__file__).parent.resolve()
TEMPLATE = Path(r"C:\Users\lizhu\OneDrive\桌面\2026.04.28_week meeting_LI ZHUOHENG.pptx")
OUTPUT = HERE / "MIPS_Assembly_Lecture.pptx"

# ---------------------------------------------------------------------------
# Style constants (BRL slate-blue scheme)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1E, 0x3A, 0x5F)       # primary - dark slate
SLATE = RGBColor(0x5C, 0x7A, 0x95)      # secondary - matches banner
LIGHT = RGBColor(0xE8, 0xEE, 0xF4)      # surface
ACCENT = RGBColor(0xED, 0x7D, 0x31)     # warm orange (theme accent2)
TEXT = RGBColor(0x21, 0x21, 0x21)       # body text
MUTED = RGBColor(0x6B, 0x72, 0x80)      # captions
CODE_BG = RGBColor(0xF6, 0xF8, 0xFA)
CODE_FG = RGBColor(0x14, 0x1B, 0x29)
COMMENT = RGBColor(0x6A, 0x73, 0x7D)
KEYWORD = RGBColor(0x9C, 0x27, 0xB0)    # purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_BODY = "Arial"
FONT_HEAD = "Arial"
FONT_CODE = "Consolas"

# Banner reserves the top ~0.8" - leave 0.95" of breathing room
TITLE_TOP = Inches(0.95)
TITLE_HEIGHT = Inches(0.7)
BODY_TOP = Inches(1.7)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
LEFT_MARGIN = Inches(0.5)
CONTENT_WIDTH = Inches(12.333)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def remove_all_slides(prs):
    """Remove every slide from the presentation in place."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_slide(prs, layout_idx=0):
    """Add a slide using the given layout (0 = with BRL banner)."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # remove any inherited placeholder text frames so we control the layout
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx in (10, 11):  # date, footer
            sp = shape._element
            sp.getparent().remove(sp)
    return slide


def add_textbox(slide, left, top, width, height, *, fill=None, line=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.line.fill.background() if line is None else None
    if line is None:
        box.line.fill.background()
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    return box, tf


def add_title(slide, text, color=NAVY, size=32):
    box, tf = add_textbox(slide, LEFT_MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT)
    tf.margin_left = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = FONT_HEAD
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return box


def add_subtitle(slide, text, top=Inches(1.55), color=MUTED, size=14):
    box, tf = add_textbox(slide, LEFT_MARGIN, top, CONTENT_WIDTH, Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.italic = True


def add_accent_bar(slide, top=Inches(1.62), color=ACCENT, width=Inches(0.6)):
    """No-op: per design rules, accent lines under titles read as AI-generated.
    Visual interest comes from the BRL banner + the navy/slate panels instead."""
    return None


def add_bullets(slide, items, *, left=None, top=None, width=None, height=None,
                font_size=18, line_spacing=1.25, bullet_color=ACCENT):
    """Add a textbox with one paragraph per item.

    `items` is a list of either:
      - "text"                 -> normal bullet
      - ("text", "sub-text")   -> bold header + indented description
      - ("HEAD::", "...")      -> use 'HEAD' as a tagged label (bold accent)
    """
    left = left or LEFT_MARGIN
    top = top or BODY_TOP
    width = width or CONTENT_WIDTH
    height = height or Inches(5.4)
    box, tf = add_textbox(slide, left, top, width, height)
    tf.margin_left = Inches(0.05)
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run()
            r1.text = "▸ "
            r1.font.color.rgb = bullet_color
            r1.font.name = FONT_BODY
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r2 = p.add_run()
            r2.text = head
            r2.font.bold = True
            r2.font.color.rgb = NAVY
            r2.font.name = FONT_BODY
            r2.font.size = Pt(font_size)
            if body:
                r3 = p.add_run()
                r3.text = "  -  " + body
                r3.font.color.rgb = TEXT
                r3.font.name = FONT_BODY
                r3.font.size = Pt(font_size)
        else:
            r1 = p.add_run()
            r1.text = "▸ "
            r1.font.color.rgb = bullet_color
            r1.font.name = FONT_BODY
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r2 = p.add_run()
            r2.text = item
            r2.font.color.rgb = TEXT
            r2.font.name = FONT_BODY
            r2.font.size = Pt(font_size)
    return box


def add_code_block(slide, code, *, left=None, top=None, width=None, height=None,
                   font_size=14):
    """Render a monospaced code block with light background and basic colour
    coding for comments (#) and labels (lines ending with :)."""
    left = left or LEFT_MARGIN
    top = top or BODY_TOP
    width = width or CONTENT_WIDTH
    if height is None:
        # estimate: ~0.22" per line at 14pt + padding
        line_count = code.count("\n") + 1
        height = Emu(min(int(Inches(0.22) * line_count) + Inches(0.3), Inches(5.6)))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.04
    box.line.color.rgb = SLATE
    box.line.width = Pt(0.5)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    tf.word_wrap = False
    lines = code.split("\n")
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        # Detect comment portion
        comment_idx = line.find("#")
        if comment_idx >= 0:
            code_part = line[:comment_idx]
            comment_part = line[comment_idx:]
        else:
            code_part = line
            comment_part = ""
        if code_part:
            r = p.add_run()
            r.text = code_part
            r.font.name = FONT_CODE
            r.font.size = Pt(font_size)
            stripped = code_part.strip()
            if stripped.endswith(":") and not stripped.startswith("#"):
                r.font.color.rgb = NAVY
                r.font.bold = True
            elif stripped.startswith(".data") or stripped.startswith(".text"):
                r.font.color.rgb = KEYWORD
                r.font.bold = True
            else:
                r.font.color.rgb = CODE_FG
        if comment_part:
            r = p.add_run()
            r.text = comment_part
            r.font.name = FONT_CODE
            r.font.size = Pt(font_size)
            r.font.color.rgb = COMMENT
            r.font.italic = True
    return box


def add_callout(slide, text, left, top, width, height, *,
                fill=NAVY, font_color=WHITE, font_size=16, bold=True, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.08
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.shadow.inherit = False
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(font_size)
    r.font.color.rgb = font_color
    r.font.bold = bold
    return box


def add_label(slide, text, left, top, width, height, *,
              color=TEXT, font_size=14, bold=False, italic=False, align=PP_ALIGN.LEFT):
    box, tf = add_textbox(slide, left, top, width, height)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.2
    for line in text.split("\n"):
        if p.runs:
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = 1.2
        r = p.add_run()
        r.text = line
        r.font.name = FONT_BODY
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return box


def add_footer(slide, text="MIPS Assembly  -  Computer Architecture, Hanyang ERICA"):
    box, tf = add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED
    r.font.italic = True


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_title(prs):
    s = add_slide(prs, 0)
    # Hero panel
    panel = slide_panel(s, NAVY, top=Inches(2.2), height=Inches(3.6))
    add_label(s,
              "MIPS Assembly Programming",
              Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.2),
              color=WHITE, font_size=44, bold=True)
    add_label(s,
              "From your first 'Hello World' to recursive functions on the stack",
              Inches(1.0), Inches(3.7), Inches(11.3), Inches(0.6),
              color=LIGHT, font_size=18, italic=True)
    add_label(s,
              "Hanyang University ERICA Campus",
              Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.45),
              color=WHITE, font_size=16, bold=True)
    add_label(s,
              "Department of Robotics  -  Computer Architecture",
              Inches(1.0), Inches(4.95), Inches(11.3), Inches(0.4),
              color=LIGHT, font_size=14)
    add_label(s,
              "Instructor: Prof. Bumjin Jang",
              Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.4),
              color=LIGHT, font_size=14)
    add_label(s,
              "Presenter: LI Zhuoheng",
              Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.4),
              color=NAVY, font_size=13, italic=True)
    add_label(s, "1-hour lecture  -  20 worked examples",
              Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.4),
              color=MUTED, font_size=12, italic=True)


def slide_panel(slide, fill, top=Inches(1.0), left=Inches(0.5), width=None, height=None):
    width = width or Inches(12.333)
    height = height or Inches(5.7)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.03
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.shadow.inherit = False
    return box


def slide_outline(prs):
    s = add_slide(prs, 0)
    add_title(s, "Today's 1-Hour Tour")
    add_accent_bar(s)
    add_subtitle(s, "Five themes, twenty programs, one live demo")
    rows = [
        ("0:00 - 0:10", "Foundations",  "MIPS architecture, MARS IDE, register conventions"),
        ("0:10 - 0:20", "Output",       "Print strings, characters, integers, floats, doubles"),
        ("0:20 - 0:35", "Arithmetic",   "add / sub / mul / mult / div with HI-LO + sll trick"),
        ("0:35 - 0:50", "Functions",    "jal / jr, calling convention, the stack, recursion"),
        ("0:50 - 0:55", "Input",        "Reading numbers and strings from the user"),
        ("0:55 - 1:00", "Demo & Q&A",   "Step through factorial(5) live in MARS"),
    ]
    y = Inches(2.1)
    row_h = Inches(0.7)
    for time, title, desc in rows:
        add_callout(s, time, Inches(0.5), y, Inches(2.1), row_h, fill=NAVY, font_size=14)
        add_callout(s, title, Inches(2.75), y, Inches(2.6), row_h, fill=SLATE, font_size=15)
        add_label(s, desc, Inches(5.5), y + Inches(0.18), Inches(7.4), Inches(0.5),
                  color=TEXT, font_size=15)
        y += Inches(0.85)
    add_footer(s)


def slide_what_is_mips(prs):
    s = add_slide(prs, 0)
    add_title(s, "What is MIPS?")
    add_accent_bar(s)
    add_subtitle(s, "A clean, classic RISC instruction set architecture")
    add_bullets(s, [
        ("RISC", "Reduced Instruction Set Computer  -  small set of fixed-length 32-bit instructions"),
        ("Load/Store architecture", "memory is accessed only via 'lw' and 'sw'; arithmetic works on registers"),
        ("32 general-purpose registers", "all the same size, no hidden side effects"),
        ("Pipeline-friendly", "the textbook example for teaching CPU design"),
        ("Why we learn it", "understand what 'high-level' code really compiles to"),
    ], top=Inches(1.95), font_size=18, line_spacing=1.4)
    # side callout
    add_callout(s, "MIPS = Microprocessor without Interlocked Pipeline Stages",
                Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7),
                fill=LIGHT, font_color=NAVY, font_size=14, bold=True)


def slide_mips_today(prs):
    s = add_slide(prs, 0)
    add_title(s, "Why MIPS Still Matters")
    add_accent_bar(s)
    add_subtitle(s, "From the PlayStation to your future embedded device")
    cards = [
        ("Education",  "Patterson & Hennessy textbook\nUsed worldwide in CS courses",  NAVY),
        ("Embedded",   "Routers (Cisco), TVs, set-top boxes,\nautomotive controllers",  SLATE),
        ("Gaming",     "Sony PlayStation 1 & 2,\nNintendo 64, PSP",                     ACCENT),
        ("Open ISA",   "Now royalty-free,\nused in academic SoC research",              RGBColor(0x4C,0xAF,0x50)),
    ]
    x = Inches(0.5)
    w = Inches(2.95)
    gap = Inches(0.13)
    y = Inches(2.1)
    h = Inches(2.6)
    for title, text, color in cards:
        # header strip
        add_callout(s, title, x, y, w, Inches(0.6), fill=color, font_size=18)
        # body card
        body = slide_panel(s, LIGHT, top=y + Inches(0.6), left=x, width=w, height=h - Inches(0.6))
        add_label(s, text, x + Inches(0.15), y + Inches(0.8), w - Inches(0.3), h - Inches(0.8),
                  color=TEXT, font_size=14)
        x += w + gap
    add_label(s,
              "Bottom line: the lessons here transfer directly to ARM, RISC-V, and any other RISC family.",
              Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
              color=NAVY, font_size=14, italic=True, bold=True)


def slide_mars(prs):
    s = add_slide(prs, 0)
    add_title(s, "MARS: Our Simulator")
    add_accent_bar(s)
    add_subtitle(s, "MIPS Assembler and Runtime Simulator (Java, free)")
    add_bullets(s, [
        ("Edit", "syntax-highlighted .asm editor, multi-file support"),
        ("Assemble (F3)", "translates your code to machine code, reports syntax errors"),
        ("Run (F5)", "execute the program; output shows in the bottom Run I/O pane"),
        ("Step (F7)", "single-step through each instruction"),
        ("Inspect", "live view of all 32 registers and main memory while you step"),
        ("Help", "built-in syscall reference (Help -> Syscalls)"),
    ], top=Inches(1.95), font_size=17, line_spacing=1.35)
    add_callout(s,
                "Download:  https://dpetersanderson.github.io/  (requires Java)",
                Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55),
                fill=NAVY, font_size=13)


def slide_mars_tour(prs):
    s = add_slide(prs, 0)
    add_title(s, "MARS Anatomy")
    add_accent_bar(s)
    add_subtitle(s, "Four panels you will use constantly")
    quads = [
        ("1. Edit / Execute Tabs", "Edit shows your source.\nExecute shows assembled machine code."),
        ("2. Code Editor",         "Type and edit your .asm file here.\nLine numbers, syntax colours."),
        ("3. Registers Panel",     "All 32 GP registers + HI / LO + PC.\nValues update as you step."),
        ("4. Run I/O",             "syscall 4 / 1 prints appear here.\nsyscall 5 / 8 reads input here."),
    ]
    cols = 2
    cw = Inches(6.05)
    ch = Inches(2.3)
    x0, y0 = Inches(0.5), Inches(2.05)
    gap = Inches(0.2)
    for i, (h, body) in enumerate(quads):
        col = i % cols
        row = i // cols
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + gap)
        slide_panel(s, LIGHT, top=y, left=x, width=cw, height=ch)
        # header bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.55))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.shadow.inherit = False
        tf = bar.text_frame
        tf.margin_left = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.name = FONT_HEAD
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = WHITE
        # body text
        add_label(s, body, x + Inches(0.2), y + Inches(0.7),
                  cw - Inches(0.4), ch - Inches(0.7),
                  color=TEXT, font_size=14)


def slide_registers(prs):
    s = add_slide(prs, 0)
    add_title(s, "General-Purpose Registers")
    add_accent_bar(s)
    add_subtitle(s, "32 registers, named by role - learn the convention, not the number")
    rows = [
        ("$zero",        "$0",      "always 0 (cannot be written)"),
        ("$at",          "$1",      "assembler-temporary  -  do not touch"),
        ("$v0 - $v1",    "$2-$3",   "function return values, also syscall code in $v0"),
        ("$a0 - $a3",    "$4-$7",   "function arguments (and syscall arguments)"),
        ("$t0 - $t9",    "$8-$15, $24-$25", "temporaries  -  caller-saved (free for any use)"),
        ("$s0 - $s7",    "$16-$23", "saved registers  -  callee MUST preserve them"),
        ("$gp / $sp / $fp", "$28-$30", "global / stack / frame pointer"),
        ("$ra",          "$31",     "return address (set automatically by jal)"),
    ]
    # table-like layout
    y = Inches(2.0)
    rh = Inches(0.55)
    add_callout(s, "Name",   Inches(0.5), y, Inches(2.4), rh, fill=NAVY,  font_size=14)
    add_callout(s, "Number", Inches(2.95), y, Inches(2.0), rh, fill=NAVY, font_size=14)
    add_callout(s, "Purpose",Inches(5.0),  y, Inches(7.83), rh, fill=NAVY,font_size=14, align=PP_ALIGN.LEFT)
    y += rh + Inches(0.05)
    alt = False
    for name, num, purpose in rows:
        row_fill = LIGHT if alt else WHITE
        alt = not alt
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(12.33), Inches(0.5))
        bg.line.fill.background()
        bg.fill.solid()
        bg.fill.fore_color.rgb = row_fill
        bg.shadow.inherit = False
        add_label(s, name,    Inches(0.6), y + Inches(0.1), Inches(2.3), Inches(0.4),
                  color=NAVY, font_size=14, bold=True)
        add_label(s, num,     Inches(2.95), y + Inches(0.1), Inches(2.0), Inches(0.4),
                  color=MUTED, font_size=14)
        add_label(s, purpose, Inches(5.0), y + Inches(0.1), Inches(7.6), Inches(0.4),
                  color=TEXT, font_size=14)
        y += Inches(0.5)


def slide_fp_registers(prs):
    s = add_slide(prs, 0)
    add_title(s, "Floating-Point Registers (Coprocessor 1)")
    add_accent_bar(s)
    add_subtitle(s, "Separate file: $f0 - $f31, used by all FP instructions")
    add_bullets(s, [
        ("32 single-precision regs", "$f0 through $f31  -  hold 32-bit floats"),
        ("Doubles use a PAIR", "$f0/$f1, $f2/$f3, ...  -  always even/odd"),
        ("Argument convention", "$f12 receives the float / double for print syscalls"),
        ("Return convention", "$f0 holds the result for read-float / read-double syscalls"),
        ("Move between FP regs", "use mov.s (single) or mov.d (double), NOT 'move'"),
        ("Load / store", "lwc1 / swc1 for floats; l.d / s.d for doubles"),
    ], top=Inches(1.95), font_size=17, line_spacing=1.4)


def slide_anatomy(prs):
    s = add_slide(prs, 0)
    add_title(s, "Anatomy of an .asm File")
    add_accent_bar(s)
    add_subtitle(s, "Two sections, one syscall pattern - every program looks like this")
    code = """.data                          # variables / constants live here
    msg: .asciiz "Hello!"      # null-terminated string

.text                          # instructions live here
    li   $v0, 4                # service code 4 = print string
    la   $a0, msg              # argument: address of the string
    syscall                    # ask the OS to do it

    li   $v0, 10               # service code 10 = exit
    syscall                    # clean termination"""
    add_code_block(s, code, top=Inches(1.95), height=Inches(3.7), font_size=18)
    # syscall pattern callout
    add_callout(s,
                "Pattern: load syscall code into $v0  ->  load arguments  ->  syscall",
                Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
                fill=NAVY, font_size=15)


def slide_print_string(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 3 - Print a String  (syscall 4)")
    add_accent_bar(s)
    code = """.data
    msg: .asciiz "Hello World!"

.text
    li   $v0, 4          # 4 = print string
    la   $a0, msg        # address of the string
    syscall

    li   $v0, 10         # 10 = exit
    syscall"""
    add_code_block(s, code, top=Inches(1.7), width=Inches(7.5), height=Inches(4.0),
                   font_size=15)
    # explanation panel on right
    panel_left = Inches(8.3)
    panel_w = Inches(4.5)
    slide_panel(s, LIGHT, top=Inches(1.7), left=panel_left, width=panel_w, height=Inches(4.0))
    add_label(s, "Key ideas", panel_left + Inches(0.2), Inches(1.85),
              panel_w - Inches(0.4), Inches(0.4),
              color=NAVY, font_size=15, bold=True)
    add_label(s,
              ".asciiz   stores text + a 0 byte at the end.\n\n"
              "li   loads an immediate (constant) into a register.\n\n"
              "la   loads the ADDRESS of a label  -  not the value.\n\n"
              "syscall hands control to the OS to actually print.",
              panel_left + Inches(0.2), Inches(2.2), panel_w - Inches(0.4), Inches(3.5),
              color=TEXT, font_size=13)
    add_callout(s, "Output:  Hello World!",
                Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
                fill=NAVY, font_size=16)


def slide_print_char_int(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lessons 4 & 5 - Print Character / Integer")
    add_accent_bar(s)
    add_subtitle(s, "Same pattern, different syscall")
    code_a = """# syscall 11 = print one character

li   $v0, 11
li   $a0, 'A'        # ASCII 65
syscall

li   $v0, 10
syscall"""
    code_b = """# syscall 1 = print signed 32-bit int

li   $v0, 1
li   $a0, 42
syscall

li   $v0, 10
syscall"""
    add_label(s, "Print a Character", Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.4),
              color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_a, left=Inches(0.5), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(3.5), font_size=15)
    add_label(s, "Print an Integer", Inches(6.83), Inches(1.95), Inches(6.0), Inches(0.4),
              color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_b, left=Inches(6.83), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(3.5), font_size=15)
    add_callout(s,
                "syscall code in $v0  +  argument in $a0  ->  syscall",
                Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
                fill=SLATE, font_size=14)


def slide_print_float_double(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lessons 6 & 7 - Print Float / Double")
    add_accent_bar(s)
    add_subtitle(s, "FP arguments live in $f12 (and $f13 for doubles)")
    code_a = """.data
    pi: .float 3.14159

.text
    li   $v0, 2          # print float
    lwc1 $f12, pi
    syscall"""
    code_b = """.data
    pi: .double 3.141592653589793

.text
    li   $v0, 3          # print double
    l.d  $f12, pi        # loads $f12 + $f13
    syscall"""
    add_label(s, "Float (32-bit, syscall 2)", Inches(0.5), Inches(1.95),
              Inches(6.0), Inches(0.4), color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_a, left=Inches(0.5), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(3.0), font_size=15)
    add_label(s, "Double (64-bit, syscall 3)", Inches(6.83), Inches(1.95),
              Inches(6.0), Inches(0.4), color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_b, left=Inches(6.83), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(3.0), font_size=15)
    add_callout(s,
                "Doubles ALWAYS occupy a pair: $f12 + $f13, $f20 + $f21, ...",
                Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.6),
                fill=NAVY, font_size=14)
    add_callout(s,
                "Use 'lwc1' for floats and 'l.d' for doubles  -  do NOT use plain 'lw'",
                Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                fill=ACCENT, font_color=WHITE, font_size=14)


def slide_add_sub(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lessons 8 & 9 - Add and Subtract")
    add_accent_bar(s)
    add_subtitle(s, "Load operands  ->  compute  ->  print")
    code = """.data
    num1: .word 30
    num2: .word 12

.text
    lw   $t0, num1            # $t0 = 30
    lw   $t1, num2            # $t1 = 12

    add  $t2, $t0, $t1        #  -> 42       (Lesson 8)
    sub  $t3, $t0, $t1        #  -> 18       (Lesson 9)

    li   $v0, 1               # print int
    move $a0, $t2             # ('move' = pseudo for 'add ..., $zero')
    syscall"""
    add_code_block(s, code, top=Inches(1.85), height=Inches(4.6), font_size=15)
    add_callout(s,
                "All R-type arithmetic has the same shape:  op  rd, rs, rt",
                Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.55),
                fill=NAVY, font_size=14)


def slide_mul_variants(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lessons 10 & 11 - Multiplication: mul vs mult")
    add_accent_bar(s)
    add_subtitle(s, "Pseudo-instruction (easy) versus the real hardware (precise)")
    code_a = """# 'mul' - PSEUDO, 3 operands
# Result goes directly into $t2.
# Convenient if the answer fits in 32 bits.

mul  $t2, $t0, $t1     # $t2 = $t0 * $t1"""
    code_b = """# 'mult' - REAL, 2 operands
# Result is 64 bits, split across HI:LO.

mult $t0, $t1
mflo $t2               # low  32 bits
mfhi $t3               # high 32 bits"""
    add_label(s, "Lesson 10  -  mul (pseudo)", Inches(0.5), Inches(1.95),
              Inches(6.0), Inches(0.4), color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_a, left=Inches(0.5), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(3.0), font_size=14)
    add_label(s, "Lesson 11  -  mult (real)", Inches(6.83), Inches(1.95),
              Inches(6.0), Inches(0.4), color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_b, left=Inches(6.83), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(3.0), font_size=14)
    add_callout(s,
                "100000 * 50000 overflows 32 bits  ->  use mult + mfhi/mflo to keep the upper word",
                Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.6),
                fill=ACCENT, font_color=WHITE, font_size=14)
    add_callout(s,
                "Rule of thumb: prototype with 'mul', switch to 'mult' when correctness matters",
                Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                fill=NAVY, font_size=14)


def slide_sll_shift(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 12 - Multiply via Shift  (sll)")
    add_accent_bar(s)
    add_subtitle(s, "When the multiplier is a power of two, shifting is much faster")
    code = """lw   $t0, num         # $t0 = 5

sll  $t1, $t0, 3      # $t1 = 5 << 3 = 5 * 2^3 = 40

# General rule:    x << n   ==   x * (2^n)"""
    add_code_block(s, code, top=Inches(1.85), height=Inches(2.4), font_size=18)
    # Visualisation
    slide_panel(s, LIGHT, top=Inches(4.5), left=Inches(0.5), width=Inches(12.33),
                height=Inches(2.5))
    add_label(s, "Bit picture", Inches(0.7), Inches(4.6), Inches(4.0), Inches(0.4),
              color=NAVY, font_size=15, bold=True)
    add_label(s,
              "5             =  0000 0000 0000 0000 0000 0000 0000 0101\n"
              "5 << 3        =  0000 0000 0000 0000 0000 0000 0010 1000   ( = 40 )\n\n"
              "Each left shift inserts a 0 on the right -> value doubles.\n"
              "3 left shifts  ->  multiplied by 2 * 2 * 2 = 8.",
              Inches(0.7), Inches(5.05), Inches(12.0), Inches(2.0),
              color=TEXT, font_size=14)


def slide_div_variants(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lessons 13 & 14 - Division")
    add_accent_bar(s)
    add_subtitle(s, "Symmetric to multiplication: pseudo gives quotient, real gives both")
    code_a = """# 'div' - PSEUDO, 3 operands
# Stores ONLY the quotient.

div  $t2, $t0, $t1
# 17 / 5 -> $t2 = 3"""
    code_b = """# 'div' - REAL, 2 operands
# LO = quotient, HI = remainder.

div  $t0, $t1
mflo $t2          # quotient  = 3
mfhi $t3          # remainder = 2"""
    add_label(s, "Lesson 13  -  pseudo (quotient only)",
              Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.4),
              color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_a, left=Inches(0.5), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(2.6), font_size=15)
    add_label(s, "Lesson 14  -  real (quotient + remainder)",
              Inches(6.83), Inches(1.95), Inches(6.0), Inches(0.4),
              color=NAVY, font_size=16, bold=True)
    add_code_block(s, code_b, left=Inches(6.83), top=Inches(2.4), width=Inches(6.0),
                   height=Inches(2.6), font_size=15)
    add_callout(s,
                "When you need BOTH quotient and remainder, the real instruction is the only choice",
                Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
                fill=NAVY, font_size=14)
    add_callout(s,
                "Mnemonic: mfLO -> quotient (Lower)   mfHI -> remainder (Higher leftover)",
                Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6),
                fill=ACCENT, font_color=WHITE, font_size=14)


def slide_hilo(prs):
    s = add_slide(prs, 0)
    add_title(s, "Why HI and LO?")
    add_accent_bar(s)
    add_subtitle(s, "Two special 32-bit registers that hold mult / div results")
    add_bullets(s, [
        ("Outside the 32 GP registers", "you cannot use them as direct operands"),
        ("mult writes both", "HI = upper 32 bits, LO = lower 32 bits of a 64-bit product"),
        ("div writes both",  "HI = remainder, LO = quotient"),
        ("Read with mfhi / mflo", "'move from HI / LO' into a normal register"),
        ("Why a separate file?", "they're wide enough for the 64-bit result without changing the ISA"),
    ], top=Inches(1.95), font_size=17, line_spacing=1.45)
    add_callout(s,
                "Forgetting mflo / mfhi is the #1 reason 'my multiply does nothing'",
                Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
                fill=ACCENT, font_color=WHITE, font_size=14)


def slide_jal_jr(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 15 - Functions with jal / jr")
    add_accent_bar(s)
    add_subtitle(s, "Two instructions are all you need to call and return")
    code = """.data
    msg: .asciiz "Hello from a function!"

.text
main:
    jal  greet               # call greet(); $ra <- ret addr

    li   $v0, 10             # exit when we come back
    syscall

greet:
    li   $v0, 4
    la   $a0, msg
    syscall
    jr   $ra                 # return to caller (main)"""
    add_code_block(s, code, top=Inches(1.85), width=Inches(7.5), height=Inches(4.6),
                   font_size=14)
    panel_left = Inches(8.3)
    panel_w = Inches(4.5)
    slide_panel(s, LIGHT, top=Inches(1.85), left=panel_left, width=panel_w, height=Inches(4.6))
    add_label(s, "Mental model", panel_left + Inches(0.2), Inches(2.0),
              panel_w - Inches(0.4), Inches(0.4),
              color=NAVY, font_size=15, bold=True)
    add_label(s,
              "jal LABEL\n"
              "  -  jump to LABEL\n"
              "  -  AND copy the address\n"
              "       of the next instruction\n"
              "       into $ra\n\n"
              "jr $ra\n"
              "  -  jump to whatever\n"
              "       address $ra holds\n\n"
              "Together: function call.",
              panel_left + Inches(0.2), Inches(2.4),
              panel_w - Inches(0.4), Inches(4.0),
              color=TEXT, font_size=13)


def slide_args_return(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 16 - Arguments and Return Values")
    add_accent_bar(s)
    add_subtitle(s, "MIPS calling convention: $a0-$a3 in, $v0-$v1 out")
    code = """.text
main:
    li   $a0, 5              # 1st argument
    li   $a1, 7              # 2nd argument
    jal  addTwo              # result lands in $v0

    move $a0, $v0            # forward to print syscall
    li   $v0, 1
    syscall

    li   $v0, 10
    syscall

addTwo:
    add  $v0, $a0, $a1       # $v0 = $a0 + $a1
    jr   $ra"""
    add_code_block(s, code, top=Inches(1.85), width=Inches(7.5), height=Inches(5.0),
                   font_size=14)
    panel_left = Inches(8.3)
    panel_w = Inches(4.5)
    slide_panel(s, LIGHT, top=Inches(1.85), left=panel_left, width=panel_w, height=Inches(5.0))
    add_label(s, "Convention", panel_left + Inches(0.2), Inches(2.0),
              panel_w - Inches(0.4), Inches(0.4),
              color=NAVY, font_size=15, bold=True)
    add_label(s,
              "Caller side\n"
              "  -  put args in $a0, $a1, ...\n"
              "  -  jal target\n"
              "  -  read result from $v0\n\n"
              "Callee side\n"
              "  -  use $a0..$a3 as inputs\n"
              "  -  put result in $v0\n"
              "  -  jr $ra to return",
              panel_left + Inches(0.2), Inches(2.4),
              panel_w - Inches(0.4), Inches(4.4),
              color=TEXT, font_size=13)


def slide_stack_concept(prs):
    s = add_slide(prs, 0)
    add_title(s, "The Stack")
    add_accent_bar(s)
    add_subtitle(s, "A LIFO scratch space, addressed through $sp")
    add_bullets(s, [
        ("$sp = stack pointer", "always points to the lowest used byte"),
        ("Stack grows DOWNWARD", "to push, SUBTRACT from $sp"),
        ("To pop, ADD back to $sp", "after reading the value"),
        ("Each slot is 4 bytes", "for word-aligned data; align to 8 for doubles"),
    ], top=Inches(1.95), font_size=17, line_spacing=1.4, height=Inches(2.6))

    # Stack visualisation on right
    base_x = Inches(8.5)
    base_y = Inches(2.0)
    slot_w = Inches(3.4)
    slot_h = Inches(0.55)
    labels = [
        ("(higher addresses)", LIGHT, MUTED, False),
        ("caller's data",      WHITE, TEXT,  False),
        ("$ra (saved)",        SLATE, WHITE, True),
        ("$a0 (saved)",        SLATE, WHITE, True),
        ("<- $sp",             ACCENT, WHITE, True),
        ("(lower addresses,",   LIGHT, MUTED, False),
        ("stack grows here)",   LIGHT, MUTED, False),
    ]
    for i, (text, fill, fc, bold) in enumerate(labels):
        y = base_y + slot_h * i
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, base_x, y, slot_w, slot_h)
        rect.line.color.rgb = SLATE
        rect.line.width = Pt(0.5)
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.shadow.inherit = False
        tf = rect.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(13)
        r.font.color.rgb = fc
        r.font.bold = bold
    add_callout(s,
                "Push:  addi $sp, $sp, -4   ;   sw $reg, 0($sp)\n"
                "Pop :  lw   $reg, 0($sp)   ;   addi $sp, $sp, 4",
                Inches(0.5), Inches(5.0), Inches(7.83), Inches(1.4),
                fill=NAVY, font_size=14, align=PP_ALIGN.LEFT)


def slide_save_to_stack(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 17 - Saving Registers to the Stack")
    add_accent_bar(s)
    add_subtitle(s, "$s0-$s7 are callee-saved: a function MUST preserve them")
    code = """clobber:
    addi $sp, $sp, -4        # allocate 4 bytes (PUSH)
    sw   $s0, 0($sp)         # save caller's $s0

    li   $s0, 999            # callee uses $s0 freely

    lw   $s0, 0($sp)         # restore caller's $s0 (POP)
    addi $sp, $sp, 4         # free the stack space
    jr   $ra"""
    add_code_block(s, code, top=Inches(1.85), height=Inches(3.4), font_size=15)
    add_callout(s,
                "If you forget to restore $s0, the CALLER will see 999 instead of 100",
                Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
                fill=ACCENT, font_color=WHITE, font_size=14)
    add_callout(s,
                "$t0-$t9 are caller-saved -> the caller must save them if they care after a call",
                Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6),
                fill=NAVY, font_size=14)


def slide_nested(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 18 (Part 1) - Why Nested Calls Need the Stack")
    add_accent_bar(s)
    add_subtitle(s, "The inner 'jal' overwrites $ra - so the OUTER caller MUST save it first")
    add_bullets(s, [
        ("Single call?", "$ra survives - nothing else writes to it"),
        ("A function calls another?", "the inner jal CLOBBERS $ra"),
        ("Solution", "outer function pushes $ra on the stack BEFORE the inner call"),
        ("Same for $a0", "if you'll need the argument again after the call, save it too"),
        ("Recursion", "is just a function calling itself - same rule applies"),
    ], top=Inches(1.95), font_size=17, line_spacing=1.45, height=Inches(3.0))
    add_callout(s,
                "Allocate ONE stack frame at the top of the function, free it before EVERY return path",
                Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                fill=NAVY, font_size=14)


def slide_recursion_walk(prs):
    s = add_slide(prs, 0)
    add_title(s, "factorial(3) - Stack Walkthrough")
    add_accent_bar(s)
    add_subtitle(s, "Each call gets its own slot for $ra and n")
    # Vertical stacked frames - left = call sequence, right = stack picture
    add_label(s, "Call sequence", Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.4),
              color=NAVY, font_size=15, bold=True)
    seq = [
        "main calls factorial(3)",
        "  factorial(3) pushes $ra,3   then calls factorial(2)",
        "    factorial(2) pushes $ra,2   then calls factorial(1)",
        "      factorial(1) returns 1   (base case)",
        "    factorial(2) pops $ra,2   returns 1 * 2 = 2",
        "  factorial(3) pops $ra,3   returns 2 * 3 = 6",
        "main receives 6 in $v0",
    ]
    add_label(s, "\n".join(seq), Inches(0.5), Inches(2.4), Inches(7.5), Inches(4.5),
              color=TEXT, font_size=14)

    # Right column: deepest stack picture
    base_x = Inches(8.5)
    base_y = Inches(2.0)
    slot_w = Inches(3.5)
    slot_h = Inches(0.5)
    add_label(s, "Stack at deepest call", base_x, Inches(1.55),
              slot_w, Inches(0.4), color=NAVY, font_size=14, bold=True)
    frames = [
        ("n = 3",         SLATE,  WHITE),
        ("$ra (main)",    NAVY,   WHITE),
        ("n = 2",         SLATE,  WHITE),
        ("$ra (fact 3)",  NAVY,   WHITE),
        ("n = 1 (base)",  ACCENT, WHITE),
        ("<-  $sp",       LIGHT,  NAVY),
    ]
    for i, (text, fill, fc) in enumerate(frames):
        y = base_y + slot_h * i
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, base_x, y, slot_w, slot_h)
        rect.line.color.rgb = SLATE
        rect.line.width = Pt(0.5)
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.shadow.inherit = False
        tf = rect.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT_CODE if "$" in text or "n =" in text else FONT_BODY
        r.font.size = Pt(13)
        r.font.color.rgb = fc
        r.font.bold = True


def slide_factorial_code(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 18 (Part 2) - Recursive Factorial Code")
    add_accent_bar(s)
    code = """factorial:
    addi $sp, $sp, -8        # reserve two slots
    sw   $ra, 0($sp)         # save return address
    sw   $a0, 4($sp)         # save n

    slti $t0, $a0, 2         # n < 2 ?
    beq  $t0, $zero, recurse

    li   $v0, 1              # base: factorial(0)=factorial(1)=1
    addi $sp, $sp, 8
    jr   $ra

recurse:
    addi $a0, $a0, -1
    jal  factorial           # $v0 = factorial(n - 1)

    lw   $a0, 4($sp)         # restore n
    lw   $ra, 0($sp)         # restore $ra
    addi $sp, $sp, 8

    mul  $v0, $v0, $a0       # $v0 = factorial(n-1) * n
    jr   $ra"""
    add_code_block(s, code, top=Inches(1.7), height=Inches(5.2), font_size=13)


def slide_input_numbers(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lessons 19 - 21 - Reading Numbers")
    add_accent_bar(s)
    add_subtitle(s, "syscall 5 / 6 / 7  =  read int / float / double")
    rows = [
        ("syscall 5",  "read integer", "result in $v0",     "move $t0, $v0"),
        ("syscall 6",  "read float",   "result in $f0",     "mov.s $f20, $f0"),
        ("syscall 7",  "read double",  "result in $f0/$f1", "mov.d $f20, $f0"),
    ]
    y = Inches(2.0)
    add_callout(s, "Code",     Inches(0.5),  y, Inches(2.0), Inches(0.55), fill=NAVY, font_size=14)
    add_callout(s, "Reads",    Inches(2.6),  y, Inches(2.6), Inches(0.55), fill=NAVY, font_size=14)
    add_callout(s, "Result lands in", Inches(5.3),  y, Inches(3.5), Inches(0.55), fill=NAVY, font_size=14)
    add_callout(s, "Save it!", Inches(8.93), y, Inches(3.9), Inches(0.55), fill=NAVY, font_size=14)
    y += Inches(0.6)
    for code, reads, lands, save in rows:
        add_callout(s, code,  Inches(0.5),  y, Inches(2.0), Inches(0.7),
                    fill=LIGHT, font_color=NAVY, font_size=14)
        add_label(s, reads,   Inches(2.7),  y + Inches(0.18), Inches(2.6), Inches(0.5),
                  color=TEXT, font_size=15)
        add_label(s, lands,   Inches(5.4),  y + Inches(0.18), Inches(3.5), Inches(0.5),
                  color=TEXT, font_size=15)
        add_label(s, save,    Inches(9.0),  y + Inches(0.18), Inches(3.9), Inches(0.5),
                  color=NAVY, font_size=15, bold=True)
        y += Inches(0.85)
    add_callout(s,
                "Always SAVE the result before the next syscall - any syscall may reuse $v0 / $f0",
                Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7),
                fill=ACCENT, font_color=WHITE, font_size=15)


def slide_input_string(prs):
    s = add_slide(prs, 0)
    add_title(s, "Lesson 22 - Reading a String  (syscall 8)")
    add_accent_bar(s)
    add_subtitle(s, "Reserve space yourself - the OS won't allocate for you")
    code = """.data
    prompt:   .asciiz "Enter your name: "
    greeting: .asciiz "Hello, "
    buffer:   .space 64           # 64 bytes of zero memory

.text
    li   $v0, 4
    la   $a0, prompt
    syscall

    li   $v0, 8                   # read string
    la   $a0, buffer              # where to put it
    li   $a1, 64                  # max bytes (incl trailing 0)
    syscall

    li   $v0, 4
    la   $a0, greeting
    syscall

    li   $v0, 4
    la   $a0, buffer              # buffer is null-terminated -> printable
    syscall"""
    add_code_block(s, code, top=Inches(1.85), height=Inches(5.0), font_size=12)


def slide_demo(prs):
    s = add_slide(prs, 0)
    add_title(s, "Live Demo - factorial(5) in MARS")
    add_accent_bar(s)
    add_subtitle(s, "Step through the recursion and watch $sp climb / fall")
    add_bullets(s, [
        ("Open the file",    "18_recursive_factorial.asm in MARS"),
        ("Assemble (F3)",    "make sure no syntax errors"),
        ("Single-step (F7)", "watch $a0 decrement on each recursive call"),
        ("Watch the stack",  "switch the right pane to 'stack' to see $ra and n piling up"),
        ("Reach base case",  "$v0 becomes 1, then climb back: 1 -> 1 -> 2 -> 6 -> 24 -> 120"),
        ("Final output",     "120 in the Run I/O panel"),
    ], top=Inches(1.95), font_size=17, line_spacing=1.4)


def slide_pitfalls(prs):
    s = add_slide(prs, 0)
    add_title(s, "Common Pitfalls (and how to avoid them)")
    add_accent_bar(s)
    add_subtitle(s, "Almost every bug lives in this short list")
    items = [
        ("Forgot to save the input",      "the next syscall will overwrite $v0  ->  copy to $t0 first"),
        ("Forgot mflo / mfhi",            "mult / div leave the result in HI:LO  -  read it!"),
        ("$ra clobbered by inner call",   "save $ra on the stack BEFORE 'jal' inside a function"),
        ("Stack not balanced",            "every push needs a matching pop on every return path"),
        ("Used 'lw' on a float",          "use 'lwc1' / 'l.d' for FP loads"),
        ("Used 'move' on FP regs",        "use 'mov.s' / 'mov.d' instead"),
        ("Off-by-one in .space buffer",   "always reserve ROOM FOR newline + null"),
    ]
    add_bullets(s, items, top=Inches(1.95), font_size=15, line_spacing=1.35)


def slide_resources(prs):
    s = add_slide(prs, 0)
    add_title(s, "Resources")
    add_accent_bar(s)
    add_subtitle(s, "Everything is open-source and on GitHub")
    cards = [
        ("Tutorial videos",
         "Amell Peralta\nMIPS Assembly Programming Simplified\n(YouTube playlist)",
         NAVY),
        ("Simulator",
         "MARS - MIPS Assembler\nand Runtime Simulator\nhttps://dpetersanderson.github.io/",
         SLATE),
        ("Source code",
         "github.com/lizhuoh9/MIPS-Assembly\n20 worked examples,\nlessons 3 - 22",
         ACCENT),
    ]
    x = Inches(0.5)
    w = Inches(4.1)
    gap = Inches(0.13)
    y = Inches(2.1)
    h = Inches(3.6)
    for title, body, color in cards:
        add_callout(s, title, x, y, w, Inches(0.7), fill=color, font_size=18)
        slide_panel(s, LIGHT, top=y + Inches(0.7), left=x, width=w, height=h - Inches(0.7))
        add_label(s, body, x + Inches(0.2), y + Inches(0.95),
                  w - Inches(0.4), h - Inches(1.0),
                  color=TEXT, font_size=14)
        x += w + gap
    add_callout(s,
                "Recommended: clone the repo, open MARS, run each lesson while you watch the video.",
                Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
                fill=NAVY, font_size=14)


def slide_thanks(prs):
    s = add_slide(prs, 0)
    slide_panel(s, NAVY, top=Inches(2.2), height=Inches(3.6))
    add_label(s,
              "Thank You",
              Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.5),
              color=WHITE, font_size=72, bold=True, align=PP_ALIGN.CENTER)
    add_label(s,
              "Questions?",
              Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.8),
              color=LIGHT, font_size=28, italic=True, align=PP_ALIGN.CENTER)
    add_label(s,
              "github.com/lizhuoh9/MIPS-Assembly",
              Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5),
              color=MUTED, font_size=14, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)

    builders = [
        slide_title,                 # 1
        slide_outline,               # 2
        slide_what_is_mips,          # 3
        slide_mips_today,            # 4
        slide_mars,                  # 5
        slide_mars_tour,             # 6
        slide_registers,             # 7
        slide_fp_registers,          # 8
        slide_anatomy,               # 9
        slide_print_string,          # 10
        slide_print_char_int,        # 11
        slide_print_float_double,    # 12
        slide_add_sub,               # 13
        slide_mul_variants,          # 14
        slide_sll_shift,             # 15
        slide_div_variants,          # 16
        slide_hilo,                  # 17
        slide_jal_jr,                # 18
        slide_args_return,           # 19
        slide_stack_concept,         # 20
        slide_save_to_stack,         # 21
        slide_nested,                # 22
        slide_recursion_walk,        # 23
        slide_factorial_code,        # 24
        slide_input_numbers,         # 25
        slide_input_string,          # 26
        slide_demo,                  # 27
        slide_pitfalls,              # 28
        slide_resources,             # 29
        slide_thanks,                # 30
    ]
    for build in builders:
        build(prs)

    prs.save(str(OUTPUT))
    print(f"OK  -  wrote {OUTPUT}  ({len(builders)} slides)")


if __name__ == "__main__":
    main()
