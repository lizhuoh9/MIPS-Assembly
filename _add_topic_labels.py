"""
Adds a small topic label to the top-left of every slide of an EXISTING
MIPS_Assembly_Lecture.pptx (does NOT touch any other content).

Run on the user's edited file in-place. Safe to re-run.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

HERE = Path(__file__).parent.resolve()
FILE = HERE / "MIPS_Assembly_Lecture.pptx"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Banner spans roughly y=0.13" -> 0.91". We center the label vertically.
LABEL_LEFT = Inches(0.3)
LABEL_TOP = Inches(0.32)
LABEL_WIDTH = Inches(7.0)
LABEL_HEIGHT = Inches(0.35)
LABEL_NAME = "topic_label"           # so we can detect & skip on re-run

TOPICS = [
    "MIPS Assembly  -  Computer Architecture, Hanyang ERICA",       # 1  Title
    "MARS Simulator",                                                # 2
    "MARS Anatomy",                                                  # 3
    "General-Purpose Registers",                                     # 4
    "Floating-Point Registers",                                      # 5
    "Anatomy of an .asm File",                                       # 6
    "Lesson 3  -  Print a String",                                   # 7
    "Lessons 4 & 5  -  Print Character / Integer",                   # 8
    "Lessons 6 & 7  -  Print Float / Double",                        # 9
    "Lessons 8 & 9  -  Add and Subtract",                            # 10
    "Lessons 10 & 11  -  mul vs mult",                               # 11
    "Lesson 12  -  Multiply via Shift (sll)",                        # 12
    "Lessons 13 & 14  -  Division",                                  # 13
    "HI and LO Registers",                                           # 14
    "Lesson 15  -  Functions with jal / jr",                         # 15
    "Lesson 16  -  Arguments and Return Values",                     # 16
    "The Stack",                                                     # 17
    "Lesson 17  -  Saving Registers to the Stack",                   # 18
    "Lesson 18 (Part 1)  -  Why Nested Calls Need the Stack",        # 19
    "factorial(3)  -  Stack Walkthrough",                            # 20
    "Lesson 18 (Part 2)  -  Recursive Factorial Code",               # 21
    "Lessons 19-21  -  Reading Numbers",                             # 22
    "Lesson 22  -  Reading a String",                                # 23
    "Resources",                                                     # 24
    "Q&A",                                                           # 25
]


def has_topic_label(slide):
    for shape in slide.shapes:
        if shape.name == LABEL_NAME:
            return True
    return False


def remove_topic_label(slide):
    """Remove any existing topic label so we can re-add with new styling."""
    removed = 0
    for shape in list(slide.shapes):
        if shape.name == LABEL_NAME:
            sp = shape._element
            sp.getparent().remove(sp)
            removed += 1
    return removed


def add_topic_label(slide, text):
    box = slide.shapes.add_textbox(LABEL_LEFT, LABEL_TOP, LABEL_WIDTH, LABEL_HEIGHT)
    box.name = LABEL_NAME
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE


def main():
    prs = Presentation(str(FILE))
    if len(prs.slides) != len(TOPICS):
        print(f"!! WARN: pptx has {len(prs.slides)} slides, "
              f"TOPICS has {len(TOPICS)} - aborting to avoid mismatch")
        return
    added = replaced = 0
    for slide, topic in zip(prs.slides, TOPICS):
        if remove_topic_label(slide):
            replaced += 1
        else:
            added += 1
        add_topic_label(slide, topic)
    prs.save(str(FILE))
    print(f"OK  -  added {added} new, replaced {replaced} existing")
    print(f"Saved {FILE}")


if __name__ == "__main__":
    main()
